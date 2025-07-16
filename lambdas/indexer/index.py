"""
send documents representing object data to elasticsearch for supported file extensions.
note: we truncate outbound documents to DOC_SIZE_LIMIT characters
(to bound memory pressure and request size to elastic)

a little knowledge on deletes and delete markers:
if bucket versioning is on:
    - `aws s3api delete-object (no --version-id)` or `aws s3 rm`
        - push a new delete marker onto the stack with a version-id
        - generate ObjectRemoved:DeleteMarkerCreated

if bucket versioning was on and is then turned off:
    - `aws s3 rm` or `aws s3api delete-object (no --version-id)`
        - replace event at top of stack
            - if a versioned delete marker, push a new one on top of it
            - if an un-versioned delete marker, replace that marker with new marker
            with version "null" (ObjectCreate will similarly replace the same with an object
            of version "null")
            - if object, destroy object
        - generate ObjectRemoved:DeleteMarkerCreated
            - problem: no way of knowing if DeleteMarkerCreated destroyed bytes
            or just created a DeleteMarker; this is usually given by the return
            value of `delete-object` but the S3 event has no knowledge of the same
    - `aws s3api delete-object --version-id VERSION`
        - destroy corresponding delete marker or object; v may be null in which
        case it will destroy the object with version null (occurs when adding
        new objects to a bucket that aws versioned but is no longer)
        - generate ObjectRemoved:Deleted

if bucket version is off and has always been off:
    - `aws s3 rm` or `aws s3api delete-object`
        - destroy object
        - generate a single ObjectRemoved:Deleted

counterintuitive things:
    - turning off versioning doesn't mean version stack can't get deeper (by at
    least 1) as indicated above in the case where a new marker is pushed onto
    the version stack
    - both creating a delete marker (soft delete) and hard deleting a delete marker
    by providing it's version-id will result in an eventType of DeleteObject
    and $.detail.responseElements.x-amz-delete-marker = true; it is therefore
    not possible to tell the difference between a new delete marker and the deletion
    of an existing one

See docs/EventBridge.md for more
"""


import datetime
import json
import os
import pathlib
import re
from os.path import split
from typing import Optional, Tuple
from urllib.parse import unquote_plus

import boto3
import botocore
import nbformat
from dateutil.tz import tzutc
from document_queue import (
    EVENT_PREFIX,
    MAX_RETRY,
    DocumentQueue,
    get_content_index_bytes,
    get_content_index_extensions,
)
from jsonschema import ValidationError, draft7_format_checker, validate
from pdfminer.high_level import extract_text as extract_pdf_text
from tenacity import (
    retry,
    retry_if_exception,
    stop_after_attempt,
    wait_exponential,
)

from quilt_shared.const import MANIFESTS_PREFIX, NAMED_PACKAGES_PREFIX
from quilt_shared.es import (
    PACKAGE_INDEX_SUFFIX,
    get_manifest_doc_id,
    get_ptr_doc_id,
    make_elastic,
    make_s3_client,
)
from t4_lambda_shared.preview import (
    ELASTIC_LIMIT_LINES,
    extract_excel,
    extract_fcs,
    extract_parquet,
    get_bytes,
    get_preview_lines,
    trim_to_bytes,
)
from t4_lambda_shared.utils import (
    get_available_memory,
    get_quilt_logger,
    separated_env_to_iter,
    set_soft_mem_limit,
)

# translate events to S3 native names
EVENTBRIDGE_TO_S3 = {
    "PutObject": EVENT_PREFIX["Created"] + "Put",
    "CopyObject": EVENT_PREFIX["Created"] + "Copy",
    "CompleteMultipartUpload": EVENT_PREFIX["Created"] + "CompleteMultipartUpload",
    # see map_event_name for complete logic
    "DeleteObject": None,
    # "DeleteObjects" is not handled since it does not contain enough information on
    # which objects where deleted
}

# ensure that we process events of known and expected shape
EVENT_SCHEMA = {
    'type': 'object',
    'properties': {
        'awsRegion': {
            'type': 'string'
        },
        'eventName': {
            'type': 'string'
        },
        'eventTime': {
            'type': 'string',
            'format': 'date-time'
        },
        's3': {
            'type': 'object',
            'properties': {
                'bucket': {
                    'type': 'object',
                    'properties': {
                        'name': {
                            'type': 'string'
                        }
                    },
                    'required': ['name'],
                    'additionalProperties': True
                },
                'object': {
                    'type': 'object',
                    'properties': {
                        'eTag': {
                            'type': 'string'
                        },
                        'isDeleteMarker': {
                            'type': 'string'
                        },
                        'key': {
                            'type': 'string'
                        },
                        'versionId': {
                            'type': 'string'
                        }
                    },
                    'required': ['key'],
                    'additionalProperties': True
                },
            },
            'required': ['bucket', 'object'],
            'additionalProperties': True
        },
    },
    'required': ['s3', 'eventName'],
    'additionalProperties': True
}
# Max number of PDF pages to extract because it can be slow
MAX_PDF_PAGES = 100
# 10 MB, see https://amzn.to/2xJpngN
NB_VERSION = 4  # default notebook version for nbformat
# currently only affects .parquet, TODO: extend to other extensions
assert 'SKIP_ROWS_EXTS' in os.environ
SKIP_ROWS_EXTS = separated_env_to_iter('SKIP_ROWS_EXTS')
TEST_EVENT = "s3:TestEvent"
MANIFEST_INDEXER_QUEUE_URL = os.environ.get("MANIFEST_INDEXER_QUEUE_URL")
assert MANIFEST_INDEXER_QUEUE_URL, "MANIFEST_INDEXER_QUEUE_URL must be set"


logger = get_quilt_logger()
s3_client = make_s3_client()
es = make_elastic(os.environ["ES_ENDPOINT"], timeout=30)
sqs = boto3.client("sqs")


def now_like_boto3():
    """ensure timezone UTC for consistency with boto3:
    Example of what boto3 returns on head_object:
        'LastModified': datetime.datetime(2019, 11, 6, 3, 1, 16, tzinfo=tzutc()),
    """
    return datetime.datetime.now(tz=tzutc())


def get_compression(ext: str):
    """return the compression type or None if not supported"""
    return "gz" if ext == ".gz" else None


def get_normalized_extensions(key) -> Tuple[str, str]:
    """standard function turning keys into a list of (possibly empty) extensions"""
    path = pathlib.PurePosixPath(key)
    try:
        ext_last = path.suffix.lower()
        ext_next_last = path.with_suffix('').suffix.lower()
    except ValueError:
        return ("", "")

    # return in left-to-right order as they occur in the key
    return (ext_next_last, ext_last)


def infer_extensions(key, exts: Tuple[str, str], compression):
    """guess extensions if possible"""
    # Handle special case of hive partitions
    # see https://www.qubole.com/blog/direct-writes-to-increase-spark-performance/
    long_ext = "".join(exts)
    # pylint: disable=too-many-boolean-expressions)
    if (
            re.fullmatch(r".c\d{3,5}", long_ext) or re.fullmatch(r".*-c\d{3,5}$", key)
            or key.endswith("_0")
            or exts[-1] == ".pq"
            or (compression and exts[0] == ".pq")
    ):
        return ".parquet"
    elif compression:
        return exts[0]

    return exts[-1]


def should_retry_exception(exception):
    """don't retry certain 40X errors"""
    if hasattr(exception, 'response'):
        error_code = exception.response.get('Error', {}).get('Code', 218)
        return error_code not in ["402", "403", "404"]
    return False


def do_index(
        s3_client,
        doc_queue: DocumentQueue,
        event_type: str,
        *,
        bucket: str,
        etag: str,
        ext: str,
        key: str,
        last_modified: str,
        text: str = '',
        size: int = 0,
        version_id: Optional[str] = None,
        s3_tags: Optional[dict] = None,
):
    """wrap dual indexing of packages and objects"""
    logger_ = get_quilt_logger()
    # index as object (always)
    logger_.debug("%s to indexing queue (%s)", key, event_type)
    doc_queue.append(
        event_type=event_type,
        bucket=bucket,
        ext=ext,
        etag=etag,
        key=key,
        last_modified=last_modified,
        size=size,
        text=text,
        version_id=version_id,
        s3_tags=s3_tags,
    )
    # maybe index as pointer
    if index_if_pointer(
        s3_client,
        doc_queue,
        bucket=bucket,
        key=key,
    ):
        logger.debug("Indexed pointer for package %s/%s", bucket, key)


def index_if_pointer(
        s3_client,
        doc_queue: DocumentQueue,
        *,
        bucket: str,
        key: str,
) -> bool:
    """index manifest pointer files as documents in ES
        Returns:
            - True if pointer to manifest (and passes to doc_queue for indexing)
            - False if not a manifest (no attempt at indexing)
    """
    logger_ = get_quilt_logger()
    pointer_prefix, pointer_file = split(key)
    handle = pointer_prefix[len(NAMED_PACKAGES_PREFIX):]
    if (
            not pointer_file
            or not pointer_prefix.startswith(NAMED_PACKAGES_PREFIX)
            or len(handle) < 3
            or '/' not in handle
    ):
        logger_.debug("Not indexing as pointer s3://%s/%s", bucket, key)
        return False
    try:
        manifest_timestamp = int(pointer_file)
        if not 1451631600 <= manifest_timestamp <= 1767250800:
            logger_.warning("Unexpected manifest timestamp s3://%s/%s", bucket, key)
            return False
    except ValueError as err:
        logger_.debug("Non-integer manifest pointer: s3://%s/%s, %s", bucket, key, err)

    try:
        resp = s3_client.get_object(
            Bucket=bucket,
            Key=key,
        )
    except (s3_client.exceptions.NoSuchKey, s3_client.exceptions.NoSuchBucket):
        logger_.info("No pointer found: s3://%s/%s. Removing.", bucket, key)
        data = None
    else:
        manifest_hash = resp["Body"].read().decode()
        # XXX: check manifest hash is valid?
        manifest_doc_id = get_manifest_doc_id(manifest_hash)
        logger_.debug("Package hash %s found in s3://%s/%s", manifest_hash, bucket, key)
        data = {
            "join_field": {"name": "ptr", "parent": manifest_doc_id},
            "routing": manifest_doc_id,
            "ptr_name": handle,
            "ptr_tag": pointer_file,
            "ptr_last_modified": resp["LastModified"],
        }
    index = bucket + PACKAGE_INDEX_SUFFIX
    _id = get_ptr_doc_id(handle, pointer_file)
    # ES index can have multiple docs with the same _id when you specify routing,
    # so we can't rely that new doc will replace the old one.
    # We have to use delete_by_query to remove the old doc from all the shards.
    es.delete_by_query(
        index=index,
        body={
            "query": {
                "bool": {
                    "filter": [{"term": {"_id": _id}}],
                    # Don't try to delete the we're trying to index.
                    # This is to avoid conflict errors that often happen during reindexing, because
                    # we get event for every version of the latest.
                    **({"must_not": [{"term": {"join_field#mnfst": manifest_doc_id}}]} if data else {}),
                }
            },
        },
    )
    if data is not None:
        doc_queue.append_document({
            "_index": bucket + PACKAGE_INDEX_SUFFIX,
            "_id": get_ptr_doc_id(handle, pointer_file),
            "_op_type": "index" if data else "delete",
            **data,
        })

    return True


def extract_pptx(fileobj, max_size: int) -> str:
    import pptx

    out = []
    prs = pptx.Presentation(fileobj)

    def iter_text_parts():
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        yield text

    for part in iter_text_parts():
        max_size -= len(part) + 1
        if max_size < 0:
            break
        out.append(part)
    return '\n'.join(out)


def maybe_get_contents(bucket, key, inferred_ext, *, etag, version_id, s3_client, size, compression=None):
    """get the byte contents of a file if it's a target for deep indexing"""
    logger_ = get_quilt_logger()
    logger_.debug(
        "Entering maybe_get_contents (could run out of mem.) %s %s %s", bucket, key, version_id
    )
    content = ""
    if inferred_ext in get_content_index_extensions(bucket_name=bucket):
        def _get_obj():
            return retry_s3(
                "get",
                bucket,
                key,
                size,
                etag=etag,
                s3_client=s3_client,
                version_id=version_id,
            )

        if inferred_ext == ".fcs":
            obj = _get_obj()
            body, info = extract_fcs(get_bytes(obj["Body"], compression), as_html=False)
            # be smart and just send column names to ES (instead of bloated full schema)
            # if this is not an HTML/catalog preview
            content = trim_to_bytes(f"{body}\n{info}", get_content_index_bytes(bucket_name=bucket))
        elif inferred_ext == ".ipynb":
            content = trim_to_bytes(
                # we have no choice but to fetch the entire notebook, because we
                # are going to parse it
                # warning: huge notebooks could spike memory here
                get_notebook_cells(
                    bucket,
                    key,
                    size,
                    compression,
                    etag=etag,
                    s3_client=s3_client,
                    version_id=version_id
                ),
                get_content_index_bytes(bucket_name=bucket),
            )
        elif inferred_ext == ".parquet":
            if size >= get_available_memory():
                print(f"{bucket}/{key} too large to deserialize; skipping contents")
                # at least index the key and other stats, but don't overrun memory
                # and fail indexing altogether
                return ""
            obj = _get_obj()
            body, info = extract_parquet(
                get_bytes(obj["Body"], compression),
                as_html=False,
                skip_rows=(inferred_ext in SKIP_ROWS_EXTS),
                max_bytes=get_content_index_bytes(bucket_name=bucket),
            )
            # be smart and just send column names to ES (instead of bloated full schema)
            # if this is not an HTML/catalog preview
            columns = ','.join(list(info['schema']['names']))
            content = trim_to_bytes(f"{columns}\n{body}", get_content_index_bytes(bucket_name=bucket))
        elif inferred_ext == ".pdf":
            obj = _get_obj()
            content = trim_to_bytes(
                extract_pdf(get_bytes(obj["Body"], compression)),
                get_content_index_bytes(bucket_name=bucket),
            )
        elif inferred_ext in (".xls", ".xlsx"):
            obj = _get_obj()
            body, _ = extract_excel(get_bytes(obj["Body"], compression), as_html=False)
            content = trim_to_bytes(
                body,
                get_content_index_bytes(bucket_name=bucket),
            )
        elif inferred_ext == ".pptx":
            obj = _get_obj()
            content = extract_pptx(get_bytes(obj["Body"], compression), get_content_index_bytes(bucket_name=bucket))
        else:
            content = get_plain_text(
                bucket,
                key,
                size,
                compression,
                etag=etag,
                s3_client=s3_client,
                version_id=version_id
            )

    return content


def extract_pdf(file_):
    """Get plain text form PDF for searchability.
    Args:
        file_ - file-like object opened in binary mode, pointing to XLS or XLSX
    Returns:
        pdf text as a string

    Warning:
        This function can be slow. The 8-page test PDF takes ~10 sec to turn into a string.
    """
    txt = extract_pdf_text(file_, maxpages=MAX_PDF_PAGES)
    # crunch down space; extract_text inserts multiple spaces
    # between words, literal newlines, etc.
    return re.sub(r"\s+", " ", txt)


def extract_text(notebook_str):
    """ Extract code and markdown
    Args:
        * nb - notebook as a string
    Returns:
        * str - select code and markdown source (and outputs)
    Pre:
        * notebook is well-formed per notebook version 4
        * "cell_type" is defined for all cells
        * "source" defined for all "code" and "markdown" cells
    Throws:
        * Anything nbformat.reads() can throw :( which is diverse and poorly
        documented, hence the `except Exception` in handler()
    Notes:
        * Deliberately decided not to index output streams and display strings
        because they were noisy and low value
        * Tested this code against ~6400 Jupyter notebooks in
        s3://alpha-quilt-storage/tree/notebook-search/
        * Might be useful to index "cell_type" : "raw" in the future
    See also:
        * Format reference https://nbformat.readthedocs.io/en/latest/format_description.html
    """
    formatted = nbformat.reads(notebook_str, as_version=NB_VERSION)
    text = []
    for cell in formatted.get("cells", []):
        if "source" in cell and cell.get("cell_type") in ("code", "markdown"):
            text.append(cell["source"])

    return "\n".join(text)


def get_notebook_cells(bucket, key, size, compression, *, etag, s3_client, version_id):
    """extract cells for ipynb notebooks for indexing"""
    text = ""
    try:
        obj = retry_s3(
            "get",
            bucket,
            key,
            size,
            etag=etag,
            s3_client=s3_client,
            version_id=version_id
        )
        data = get_bytes(obj["Body"], compression)
        notebook = data.getvalue().decode("utf-8")
        try:
            text = extract_text(notebook)
        except (json.JSONDecodeError, nbformat.reader.NotJSONError):
            print(f"Invalid JSON in {key}.")
        except (KeyError, AttributeError) as err:
            print(f"Missing key in {key}: {err}")
        # there might be more errors than covered by test_read_notebook
        # better not to fail altogether
        except Exception as exc:  # pylint: disable=broad-except
            print(f"Exception in file {key}: {exc}")
    except UnicodeDecodeError as uni:
        print(f"Unicode decode error in {key}: {uni}")

    return text


def get_plain_text(
        bucket,
        key,
        size,
        compression,
        *,
        etag,
        s3_client,
        version_id
) -> str:
    """get plain text object contents"""
    text = ""
    try:
        obj = retry_s3(
            "get",
            bucket,
            key,
            size,
            etag=etag,
            s3_client=s3_client,
            limit=get_content_index_bytes(bucket_name=bucket),
            version_id=version_id
        )
        lines = get_preview_lines(
            obj["Body"],
            compression,
            ELASTIC_LIMIT_LINES,
            get_content_index_bytes(bucket_name=bucket),
        )
        text = '\n'.join(lines)
    except UnicodeDecodeError as ex:
        print(f"Unicode decode error in {key}", ex)

    return text


def map_event_name(event: dict):
    """transform eventbridge names into S3-like ones"""
    input_ = event["eventName"]
    if input_ in EVENTBRIDGE_TO_S3:
        if input_ == "DeleteObject":
            if event["s3"]["object"].get("isDeleteMarker"):
                return EVENT_PREFIX["Removed"] + "DeleteMarkerCreated"
            return EVENT_PREFIX["Removed"] + "Delete"
        # all non-delete events just use the map
        return EVENTBRIDGE_TO_S3[input_]
    # leave event type unchanged if we don't recognize it
    return input_


def shape_event(event: dict):
    """check event schema, return None if schema check fails"""
    logger_ = get_quilt_logger()

    try:
        validate(
            instance=event,
            schema=EVENT_SCHEMA,
            # format_checker= required for for format:date-time validation
            # (we also need strict-rfc3339 in requirements.txt)
            format_checker=draft7_format_checker,
        )
    except ValidationError as error:
        logger_.error("Invalid event format: %s\n%s", error, event)
        return None
    # be a good citizen and don't modify params
    return {
        **event,
        'eventName': map_event_name(event),
    }


def handler(event, context):
    """enumerate S3 keys in event, extract relevant data, queue events, send to
    elastic via bulk() API
    """
    logger_ = get_quilt_logger()
    # message is a proper SQS message, which either contains a single event
    # (from the bucket notification system) or batch-many events as determined
    # by enterprise/**/bulk_loader.py
    # An exception that we'll want to re-raise after the batch sends
    # TODO: handle s3:ObjectTagging:* events to keep s3_tags updated
    content_exception = None
    batch_processor = DocumentQueue(context)
    s3_client = make_s3_client()
    for message in event["Records"]:
        body = json.loads(message["body"])
        body_message = json.loads(body["Message"])
        if "Records" not in body_message:
            # could be TEST_EVENT, or another unexpected event; skip it
            logger_.error("No 'Records' key in message['body']: %s", message)
            continue
        events = body_message["Records"]
        # event is a single S3 event
        for event_ in events:
            validated = shape_event(event_)
            if not validated:
                logger_.debug("Skipping invalid event %s", event_)
                continue
            event_ = validated
            logger_.debug("Processing %s", event_)

            # XXX: this is a workaround until bulk scanner is updated to send events to EventBridge
            if event_["s3"]["object"]["key"].startswith(MANIFESTS_PREFIX):
                logger_.debug("Manifest event, sending to indexer queue: %s", event_)
                sqs.send_message(
                    QueueUrl=MANIFEST_INDEXER_QUEUE_URL,
                    # Partly mimic the structure EventBridge event
                    MessageBody=json.dumps(
                        {
                            "detail": event_,
                        },
                    ),
                )

            try:
                event_name = event_["eventName"]
                # Process all Create:* and Remove:* events
                if not any(event_name.startswith(n) for n in EVENT_PREFIX.values()):
                    logger_.warning("Skipping unknown event type: %s", event_name)
                    continue
                bucket = event_["s3"]["bucket"]["name"]
                # In the grand tradition of IE6, S3 events turn spaces into '+'
                # TODO: check if eventbridge events do the same thing with +
                key = unquote_plus(event_["s3"]["object"]["key"])
                version_id = event_["s3"]["object"].get("versionId", None)
                # ObjectRemoved:Delete does not include "eTag"
                etag = event_["s3"]["object"].get("eTag", "")
                # synthetic events from bulk scanner might define lastModified
                last_modified = (
                    event_["s3"]["object"].get("lastModified") or event_["eventTime"]
                )
                # Get two levels of extensions to handle files like .csv.gz
                ext_next_last, ext_last = get_normalized_extensions(key)
                compression = get_compression(ext_last)
                ext = ext_next_last + ext_last
                # Handle delete and deletemarker first and then continue so that
                # head_object and get_object (below) don't fail
                if event_name.startswith(EVENT_PREFIX["Removed"]):
                    do_index(
                        s3_client,
                        batch_processor,
                        event_name,
                        bucket=bucket,
                        etag=etag,
                        ext=ext,
                        key=key,
                        last_modified=last_modified,
                        version_id=version_id
                    )
                    continue
                try:
                    head = retry_s3(
                        "head",
                        bucket,
                        key,
                        s3_client=s3_client,
                        version_id=version_id,
                        etag=etag
                    )
                except botocore.exceptions.ClientError as first:
                    logger_.warning("head_object error: %s", first)
                    # "null" version sometimes results in 403s for buckets
                    # that have changed versioning, retry without it
                    if (first.response.get('Error', {}).get('Code') == "403"
                            and version_id == "null"):
                        try:
                            head = retry_s3(
                                "head",
                                bucket,
                                key,
                                s3_client=s3_client,
                                version_id=None,
                                etag=etag
                            )
                        except botocore.exceptions.ClientError as second:
                            # this will bypass the DLQ but that's the right thing to do
                            # as some listed objects may NEVER succeed head requests
                            # (e.g. foreign owner) and there's no reason to torpedo
                            # the whole batch (which might include good files)
                            logger_.warning("Retried head_object error: %s", second)
                    logger_.error("Fatal head_object, skipping event: %s", event_)
                    continue
                # backfill fields based on the head_object
                size = head["ContentLength"]
                last_modified = last_modified or head["LastModified"].isoformat()
                etag = head.get("etag") or etag
                version_id = head.get("VersionId") or version_id
                try:
                    with set_soft_mem_limit(context):
                        text = maybe_get_contents(
                            bucket,
                            key,
                            infer_extensions(key, (ext_next_last, ext_last), compression),
                            etag=etag,
                            version_id=version_id,
                            s3_client=s3_client,
                            size=size,
                            compression=compression
                        )
                # we still want an entry for this document in elastic so that, e.g.,
                # the file counts from elastic are correct
                # these exceptions can happen for a variety of reasons (e.g. glacier
                # storage class, index event arrives after delete has occurred, etc.)
                # given how common they are, we shouldn't fail the batch for this
                except Exception as exc:  # pylint: disable=broad-except
                    text = ""
                    logger_.exception("Content extraction failed %s %s", bucket, key)

                # XXX: we could replace head_object() call above with get_object(Range='bytes=0-0')
                #      which returns TagsCount, so we could optimize out get_object_tagging() call
                #      for objects without tags.

                do_index(
                    s3_client,
                    batch_processor,
                    event_name,
                    bucket=bucket,
                    etag=etag,
                    ext=ext,
                    key=key,
                    last_modified=last_modified,
                    size=size,
                    text=text,
                    version_id=version_id,
                    s3_tags=get_object_tagging(
                        s3_client=s3_client,
                        bucket=bucket,
                        key=key,
                        version_id=version_id,
                    ),
                )

            except botocore.exceptions.ClientError as boto_exc:
                if not should_retry_exception(boto_exc):
                    logger_.warning("Skipping non-fatal exception: %s", boto_exc)
                    continue
                logger_.critical("Failed record: %s, %s", event, boto_exc)
                raise boto_exc
    # flush the queue
    batch_processor.send_all()


def retry_s3(
        operation,
        bucket,
        key,
        size=None,
        limit=None,
        *,
        etag,
        version_id,
        s3_client
):
    """retry head or get operation to S3 with; stop before we run out of time.
    retry is necessary since, due to eventual consistency, we may not
    always get the required version of the object.
    """
    logger_ = get_quilt_logger()

    if operation == "head":
        function_ = s3_client.head_object
    elif operation == "get":
        function_ = s3_client.get_object
    else:
        raise ValueError(f"unexpected operation: {operation}")
    # Keyword arguments to function_
    arguments = {
        "Bucket": bucket,
        "Key": key
    }
    if operation == 'get' and size and limit:
        # can only request range if file is not empty
        arguments['Range'] = f"bytes=0-{min(size, limit) - 1}"
    if version_id:
        arguments['VersionId'] = version_id
    elif etag:
        arguments['IfMatch'] = etag

    logger_.debug("Entering @retry: %s, %s", operation, arguments)

    @retry(
        # debug
        reraise=True,
        stop=stop_after_attempt(MAX_RETRY),
        wait=wait_exponential(multiplier=2, min=4, max=10),
        retry=(retry_if_exception(should_retry_exception))
    )
    def call():
        """local function so we can set stop_after_delay dynamically"""
        # TODO: remove all this, stop_after_delay is not dynamically loaded anymore
        return function_(**arguments)

    return call()


def get_object_tagging(*, s3_client, bucket: str, key: str, version_id: Optional[str]) -> Optional[dict]:
    params = {
        "Bucket": bucket,
        "Key": key,
    }
    if version_id:
        params["VersionId"] = version_id

    try:
        s3_tags = s3_client.get_object_tagging(**params)["TagSet"]
        return {t["Key"]: t["Value"] for t in s3_tags}
    except botocore.exceptions.ClientError as e:
        if e.response["Error"]["Code"] != "AccessDenied":
            raise
        get_quilt_logger().error(
            "AccessDenied while getting tags for Bucket=%s, Key=%s, VersionId=%s",
            bucket, key, version_id
        )
        return None
